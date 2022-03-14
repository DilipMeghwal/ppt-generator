from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
# import pandas as pd
import openpyxl
from openpyxl_image_loader import SheetImageLoader

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)

# to set the slide background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

# set title text
title = slide.placeholders[0]
run = title.text_frame.paragraphs[0].add_run()
font = run.font
font.name = 'Calibri'
font.size = Pt(66)
font.bold = True
run.text = "CAVIER"
run.font.color.rgb = RGBColor(255, 255, 255)
# title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

p = title.text_frame.add_paragraph()
run = p.add_run()
font = run.font
font.name = 'Calibri'
font.size = Pt(54)
# font.bold = True
run.text = "SOCIAL"
run.font.color.rgb = RGBColor(255, 255, 255)

# set subtitle text
subtitle = slide.placeholders[1]
run = subtitle.text_frame.paragraphs[0].add_run()
font = run.font
font.name = 'Calibri'
font.size = Pt(48)
font.bold = True
run.text = "COMPETITOR REPORT"
run.font.color.rgb = RGBColor(255, 255, 255)
# subtitle.text = "COMPETITOR REPORT"
# subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

p = subtitle.text_frame.add_paragraph()
run = p.add_run()
font = run.font
font.name = 'Calibri'
font.size = Pt(32)
# font.bold = True
run.text = "Prepared for"
run.font.color.rgb = RGBColor(255, 255, 255)

# -------------------- add new slide ----------------------
# excel_data_df = pd.read_excel('data.xlsx', sheet_name='Set_01')
# # print whole sheet data
# print(excel_data_df)
#
# for row in excel_data_df.itertuples():
#     print(row[0])
#loading the Excel File and the sheet
pxl_doc = openpyxl.load_workbook('data.xlsx')
sheet = pxl_doc['Set_01']

#calling the image_loader
image_loader = SheetImageLoader(sheet)

for x in range(2, sheet.max_column):
    finalString = ""
    for y in range(x, sheet.max_row + 1):
        print(sheet.cell(row=x, column=y).value)
        finalString = finalString + str(sheet.cell(row=x, column=y).value) + '\n'
    print(finalString)
    title_slide_layout_template = prs.slide_layouts[6]
    slide_template = prs.slides.add_slide(title_slide_layout_template)
    left = top = width = height = Inches(1.0)
    # to set the slide background color
    background = slide_template.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    textbox_username = slide_template.shapes.add_textbox(Inches(3), Inches(1.5), Inches(3), Inches(1))
    p = textbox_username.text_frame.add_paragraph()
    p.alignment = PP_ALIGN.JUSTIFY
    run = p.add_run()
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(32)
    # font.bold = True
    run.text = "@" + finalString
    run.font.color.rgb = RGBColor(255, 255, 255)
    #
    # # textbox_followers = slide_template.shapes.add_textbox(Inches(3), Inches(1.5), Inches(3), Inches(1))
    # # p = textbox_followers.text_frame.add_paragraph()
    # # p.alignment = PP_ALIGN.JUSTIFY
    # run = p.add_run()
    # font = run.font
    # font.name = 'Calibri'
    # font.size = Pt(32)
    # # font.bold = True
    # run.text = str(sheet.cell(row=x, column=2).value) + " followers"
    # run.font.color.rgb = RGBColor(255, 255, 255)
    #
    # # image = image_loader.get('E5')
    # # image.show()
    # # pic = slide_template.shapes.add_picture(image, left, top, width, height)
    # # pic.left = int((int(prs.slide_width) - int(pic.width)) / 2)
    # # pic.auto_shape_type = MSO_AUTO_SHAPE_TYPE.OVAL

prs.save('test.pptx')
